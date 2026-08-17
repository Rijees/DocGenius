from pathlib import Path
import io
import os
import re
import shutil
import zipfile

import fitz
from pypdf import PdfReader, PdfWriter
from PIL import Image, ImageOps
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4

TOOL_REGISTRY = {
    "merge-pdf": {"name":"Merge PDF","icon":"⇄","category":"Organize PDF","description":"Combine multiple PDF files into one document.","accept":".pdf","multiple":True},
    "split-pdf": {"name":"Split PDF","icon":"✂","category":"Organize PDF","description":"Split a PDF into separate page files.","accept":".pdf"},
    "extract-pages": {"name":"Extract Pages","icon":"▣","category":"Organize PDF","description":"Extract selected pages into a new PDF.","accept":".pdf"},
    "rotate-pdf": {"name":"Rotate PDF","icon":"↻","category":"Organize PDF","description":"Rotate all pages by 90, 180 or 270 degrees.","accept":".pdf"},
    "reorder-pages": {"name":"Reorder Pages","icon":"☷","category":"Organize PDF","description":"Rearrange page order.","accept":".pdf"},
    "delete-pages": {"name":"Delete Pages","icon":"⌫","category":"Organize PDF","description":"Remove selected pages from a PDF.","accept":".pdf"},
    "duplicate-pages": {"name":"Duplicate Pages","icon":"＋","category":"Organize PDF","description":"Duplicate selected pages.","accept":".pdf"},
    "compress-pdf": {"name":"Compress PDF","icon":"◉","category":"Optimize PDF","description":"Reduce PDF file size.","accept":".pdf"},
    "pdf-to-jpg": {"name":"PDF to JPG","icon":"▧","category":"Convert from PDF","description":"Convert PDF pages to JPG images.","accept":".pdf"},
    "jpg-to-pdf": {"name":"JPG to PDF","icon":"▤","category":"Convert to PDF","description":"Create a PDF from JPG images.","accept":".jpg,.jpeg,.png","multiple":True},
    "pdf-to-png": {"name":"PDF to PNG","icon":"▧","category":"Convert from PDF","description":"Convert PDF pages to PNG images.","accept":".pdf"},
    "png-to-pdf": {"name":"PNG to PDF","icon":"▤","category":"Convert to PDF","description":"Create a PDF from PNG images.","accept":".png,.jpg,.jpeg","multiple":True},
    "pdf-to-word": {"name":"PDF to Word","icon":"W","category":"Convert from PDF","description":"Convert PDF text into a DOCX document.","accept":".pdf"},
    "word-to-pdf": {"name":"Word to PDF","icon":"W","category":"Convert to PDF","description":"Convert a DOCX file to PDF when Office/LibreOffice is available.","accept":".docx"},
    "pdf-to-excel": {"name":"PDF to Excel","icon":"X","category":"Convert from PDF","description":"Extract detected tables into an XLSX workbook.","accept":".pdf"},
    "excel-to-pdf": {"name":"Excel to PDF","icon":"X","category":"Convert to PDF","description":"Convert an XLSX file to PDF via Office/LibreOffice.","accept":".xlsx"},
    "pdf-to-powerpoint": {"name":"PDF to PowerPoint","icon":"P","category":"Convert from PDF","description":"Create a PPTX with PDF page images.","accept":".pdf"},
    "powerpoint-to-pdf": {"name":"PowerPoint to PDF","icon":"P","category":"Convert to PDF","description":"Convert PPTX to PDF via Office/LibreOffice.","accept":".pptx"},
    "html-to-pdf": {"name":"HTML to PDF","icon":"H","category":"Convert to PDF","description":"Turn HTML content into a PDF.","accept":".html,.htm"},
    "txt-to-pdf": {"name":"TXT to PDF","icon":"T","category":"Convert to PDF","description":"Convert plain text into a PDF.","accept":".txt"},
    "protect-pdf": {"name":"Protect PDF","icon":"⌑","category":"PDF Security","description":"Encrypt a PDF with a password.","accept":".pdf"},
    "unlock-pdf": {"name":"Unlock PDF","icon":"⌑","category":"PDF Security","description":"Remove password protection when the supplied password is valid.","accept":".pdf"},
    "watermark-pdf": {"name":"Watermark PDF","icon":"W","category":"Edit PDF","description":"Add a text watermark to every page.","accept":".pdf"},
    "page-numbers": {"name":"Page Numbers","icon":"#","category":"Edit PDF","description":"Add page numbers to every page.","accept":".pdf"},
    "header-footer": {"name":"Header & Footer","icon":"↕","category":"Edit PDF","description":"Add simple headers and footers.","accept":".pdf"},
    "crop-pdf": {"name":"Crop PDF","icon":"□","category":"Edit PDF","description":"Crop page margins.","accept":".pdf"},
    "grayscale-pdf": {"name":"Grayscale PDF","icon":"◐","category":"Edit PDF","description":"Convert pages to grayscale.","accept":".pdf"},
    "pdf-metadata": {"name":"PDF Metadata","icon":"ⓘ","category":"PDF Utilities","description":"View or update basic PDF metadata.","accept":".pdf"},
    "pdf-repair": {"name":"PDF Repair","icon":"🔧","category":"PDF Utilities","description":"Rewrite readable PDF pages into a fresh PDF.","accept":".pdf"},
    "pdf-booklet": {"name":"PDF Booklet","icon":"▥","category":"PDF Utilities","description":"Prepare a simple booklet page order.","accept":".pdf"},
    "pdf-poster": {"name":"PDF Poster","icon":"▦","category":"PDF Utilities","description":"Scale a PDF page for poster-style output.","accept":".pdf"},
    "pdf-to-text": {"name":"PDF to Text","icon":"T","category":"Extract Data","description":"Extract text from PDF pages.","accept":".pdf"},
    "text-to-pdf": {"name":"Text to PDF","icon":"T","category":"Convert to PDF","description":"Create a PDF from entered text.","accept":".txt"},
    "extract-images": {"name":"PDF Images Extractor","icon":"▧","category":"Extract Data","description":"Extract embedded images from a PDF.","accept":".pdf"},
    "pdf-info": {"name":"PDF Information","icon":"ⓘ","category":"PDF Utilities","description":"Inspect pages, size, metadata and dimensions.","accept":".pdf"},
    "flatten-pdf": {"name":"Flatten PDF","icon":"▰","category":"PDF Utilities","description":"Render pages into a flattened PDF.","accept":".pdf"},
    "add-blank-pages": {"name":"Add Blank Pages","icon":"＋","category":"Organize PDF","description":"Append blank pages to a PDF.","accept":".pdf"},
    "remove-blank-pages": {"name":"Remove Blank Pages","icon":"⌫","category":"Organize PDF","description":"Remove pages with no meaningful text or images.","accept":".pdf"},
    "pdf-compare": {"name":"PDF Compare","icon":"⇆","category":"Advanced","description":"Compare text content from two PDFs.","accept":".pdf","multiple":True},
    "pdf-chat": {"name":"PDF Chat","icon":"✦","category":"AI Features","description":"AI-ready document chat interface.","accept":".pdf"},
    "ai-summary": {"name":"AI Document Summary","icon":"✦","category":"AI Features","description":"AI-ready document summarization interface.","accept":".pdf"},
    "ocr-pdf": {"name":"OCR PDF","icon":"⌕","category":"AI Features","description":"OCR extension point for scanned PDFs.","accept":".pdf"},
    "esign-pdf": {"name":"E-Sign PDF","icon":"✎","category":"Future Tools","description":"Signature workflow extension point.","accept":".pdf"},
    "scanner": {"name":"Document Scanner","icon":"▣","category":"Future Tools","description":"Camera/scanner workflow extension point.","accept":".jpg,.jpeg,.png","multiple":True},
}

def outpath(folder, name):
    p = Path(folder) / name
    p.parent.mkdir(parents=True, exist_ok=True)
    return str(p)

def pages_arg(options, default="1"):
    return [int(x)-1 for x in re.split(r"[\s,]+", options.get("pages", default).strip()) if x.isdigit()]

def merge(paths, output):
    writer = PdfWriter()
    for p in paths:
        for page in PdfReader(p).pages:
            writer.add_page(page)
    with open(output, "wb") as f: writer.write(f)

def split(path, output_dir):
    reader = PdfReader(path)
    zip_path = Path(output_dir) / "docgenius_split_pages.zip"
    temp = Path(output_dir) / "split_pages"
    if temp.exists(): shutil.rmtree(temp)
    temp.mkdir(parents=True)
    names=[]
    for i, page in enumerate(reader.pages, 1):
        w=PdfWriter(); w.add_page(page)
        fp=temp/f"page_{i}.pdf"
        with open(fp,"wb") as f:w.write(f)
        names.append(fp)
    with zipfile.ZipFile(zip_path,"w",zipfile.ZIP_DEFLATED) as z:
        for fp in names:z.write(fp,fp.name)
    return str(zip_path)

def images_to_pdf(paths, output):
    imgs=[]
    for p in paths:
        im=Image.open(p).convert("RGB")
        imgs.append(im)
    if not imgs: raise ValueError("No images")
    imgs[0].save(output, save_all=True, append_images=imgs[1:])

def pdf_to_images(path, output_dir, ext):
    doc=fitz.open(path)
    temp=Path(output_dir)/f"pdf_pages_{ext}"
    if temp.exists(): shutil.rmtree(temp)
    temp.mkdir(parents=True)
    for i,page in enumerate(doc,1):
        pix=page.get_pixmap(matrix=fitz.Matrix(1.7,1.7), alpha=False)
        pix.save(str(temp/f"page_{i}.{ext}"))
    zip_path=Path(output_dir)/f"docgenius_pdf_to_{ext}.zip"
    with zipfile.ZipFile(zip_path,"w",zipfile.ZIP_DEFLATED) as z:
        for fp in sorted(temp.iterdir()): z.write(fp,fp.name)
    return str(zip_path)

def rotate(path, output, angle):
    doc=fitz.open(path)
    for page in doc: page.set_rotation(angle % 360)
    doc.save(output)
    return output

def delete_pages(path, output, indexes):
    reader=PdfReader(path); writer=PdfWriter()
    for i,p in enumerate(reader.pages):
        if i not in indexes: writer.add_page(p)
    with open(output,"wb") as f: writer.write(f)
    return output

def extract_pages(path, output, indexes):
    reader=PdfReader(path); writer=PdfWriter()
    for i in indexes:
        if 0 <= i < len(reader.pages): writer.add_page(reader.pages[i])
    with open(output,"wb") as f: writer.write(f)
    return output

def protect(path, output, password):
    r=PdfReader(path); w=PdfWriter()
    for p in r.pages:w.add_page(p)
    w.encrypt(password)
    with open(output,"wb") as f:w.write(f)
    return output

def unlock(path, output, password):
    r=PdfReader(path)
    if r.is_encrypted:
        if not r.decrypt(password): raise ValueError("Incorrect PDF password")
    w=PdfWriter()
    for p in r.pages:w.add_page(p)
    with open(output,"wb") as f:w.write(f)
    return output

def watermark(path, output, text):
    src=fitz.open(path)
    for page in src:
        rect=page.rect
        page.insert_text((rect.width/2-100, rect.height/2), text or "DocGenius", fontsize=28, rotate=45, color=(0.4,0.4,0.4), fill_opacity=0.22)
    src.save(output)
    return output

def add_numbers(path, output):
    doc=fitz.open(path)
    for i,page in enumerate(doc,1):
        page.insert_text((page.rect.width/2-12,page.rect.height-20), str(i), fontsize=10, color=(0.25,0.25,0.25))
    doc.save(output); return output

def grayscale(path, output):
    doc=fitz.open(path)
    for page in doc:
        pix=page.get_pixmap(matrix=fitz.Matrix(1,1), alpha=False, colorspace=fitz.csGRAY)
        page.clean_contents()
        page.insert_image(page.rect, pixmap=pix)
    doc.save(output); return output

def pdf_to_text(path, output):
    doc=fitz.open(path)
    with open(output,"w",encoding="utf-8") as f:
        for i,p in enumerate(doc,1):
            f.write(f"\n--- Page {i} ---\n{p.get_text()}\n")
    return output

def pdf_info(path, output):
    doc=fitz.open(path)
    info={"pages":doc.page_count,"metadata":doc.metadata,"file_size_bytes":Path(path).stat().st_size}
    Path(output).write_text(__import__("json").dumps(info,indent=2),encoding="utf-8")
    return output

def repair(path, output):
    doc=fitz.open(path); new=fitz.open()
    new.insert_pdf(doc); new.save(output); return output

def add_blank(path, output, count):
    r=PdfReader(path); w=PdfWriter()
    for p in r.pages:w.add_page(p)
    for _ in range(count): w.add_blank_page()
    with open(output,"wb") as f:w.write(f)
    return output

def remove_blank(path, output):
    doc=fitz.open(path); new=fitz.open()
    for p in doc:
        text=p.get_text().strip()
        imgs=len(p.get_images(full=True))
        if text or imgs: new.insert_pdf(doc, from_page=p.number, to_page=p.number)
    new.save(output); return output

def pdf_to_pptx(path, output):
    from pptx import Presentation
    from pptx.util import Inches
    doc=fitz.open(path); prs=Presentation()
    prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    for i,p in enumerate(doc):
        pix=p.get_pixmap(matrix=fitz.Matrix(1.3,1.3), alpha=False)
        img=io.BytesIO(pix.tobytes("png"))
        slide=prs.slides[0] if i==0 and len(prs.slides)==1 and len(prs.slides[0].shapes)==0 else prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img,0,0,width=prs.slide_width,height=prs.slide_height)
    prs.save(output); return output

def pdf_to_word(path, output):
    from docx import Document
    docx=Document(); doc=fitz.open(path)
    for i,p in enumerate(doc):
        if i: docx.add_page_break()
        docx.add_paragraph(p.get_text())
    docx.save(output); return output

def txt_to_pdf(path, output):
    c=canvas.Canvas(output,pagesize=A4); w,h=A4
    y=h-50
    text=Path(path).read_text(encoding="utf-8",errors="ignore")
    for line in text.splitlines():
        if y<50: c.showPage(); y=h-50
        c.drawString(45,y,line[:110]); y-=14
    c.save(); return output

def text_to_pdf_from_option(text, output):
    c=canvas.Canvas(output,pagesize=A4); w,h=A4
    y=h-50
    for line in text.splitlines():
        if y<50:c.showPage();y=h-50
        c.drawString(45,y,line[:110]);y-=14
    c.save();return output

def run_tool(slug, paths, options, output_dir):
    out = lambda name: outpath(output_dir, name)
    p = paths[0] if paths else None

    if slug=="merge-pdf":
        o=out("docgenius_merged.pdf"); merge(paths,o); return o
    if slug=="split-pdf": return split(p, output_dir)
    if slug=="extract-pages":
        o=out("docgenius_extracted.pdf"); return extract_pages(p,o,pages_arg(options))
    if slug=="delete-pages":
        o=out("docgenius_deleted.pdf"); return delete_pages(p,o,pages_arg(options))
    if slug=="rotate-pdf":
        o=out("docgenius_rotated.pdf"); return rotate(p,o,int(options.get("angle","90")))
    if slug=="compress-pdf":
        o=out("docgenius_compressed.pdf"); doc=fitz.open(p); doc.save(o,garbage=4,deflate=True,clean=True); return o
    if slug=="pdf-to-jpg": return pdf_to_images(p,output_dir,"jpg")
    if slug=="pdf-to-png": return pdf_to_images(p,output_dir,"png")
    if slug in ("jpg-to-pdf","png-to-pdf"):
        o=out("docgenius_images.pdf"); images_to_pdf(paths,o); return o
    if slug=="pdf-to-word":
        o=out("docgenius_document.docx"); return pdf_to_word(p,o)
    if slug=="pdf-to-powerpoint":
        o=out("docgenius_presentation.pptx"); return pdf_to_pptx(p,o)
    if slug=="protect-pdf":
        o=out("docgenius_protected.pdf"); return protect(p,o,options.get("password",""))
    if slug=="unlock-pdf":
        o=out("docgenius_unlocked.pdf"); return unlock(p,o,options.get("password",""))
    if slug=="watermark-pdf":
        o=out("docgenius_watermarked.pdf"); return watermark(p,o,options.get("text","DocGenius"))
    if slug=="page-numbers":
        o=out("docgenius_numbered.pdf"); return add_numbers(p,o)
    if slug=="grayscale-pdf":
        o=out("docgenius_grayscale.pdf"); return grayscale(p,o)
    if slug=="pdf-to-text":
        o=out("docgenius_text.txt"); return pdf_to_text(p,o)
    if slug=="pdf-info":
        o=out("docgenius_info.json"); return pdf_info(p,o)
    if slug=="pdf-repair":
        o=out("docgenius_repaired.pdf"); return repair(p,o)
    if slug=="add-blank-pages":
        o=out("docgenius_blank_pages.pdf"); return add_blank(p,o,int(options.get("count","1")))
    if slug=="remove-blank-pages":
        o=out("docgenius_no_blank_pages.pdf"); return remove_blank(p,o)
    if slug=="txt-to-pdf":
        o=out("docgenius_text.pdf"); return txt_to_pdf(p,o)
    if slug=="text-to-pdf":
        o=out("docgenius_text.pdf"); return text_to_pdf_from_option(options.get("text",""),o)
    # Features that are intentionally extension points or need external Office/AI/OCR integration.
    return None
